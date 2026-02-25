#!/usr/bin/env python3
"""Build a PPTX from the Metals Class Reveal.js presentation.

18 slides:
  1  Title
  2  What About Metal?
  3  Ferrous vs. Non-Ferrous
  4  Why Do You Care?
  5  Metal Identification Methods
  6  Mild Steel (1018 / A36)
  7  4140 Chrome-Moly Steel
  8  Tool Steel (A2 / D2)
  9  Cast Iron
 10  6061 Aluminum
 11  Copper
 12  Brass
 13  Titanium
 14  Stainless Steel (304 / 316)
 15  MIG Welding
 16  Lathe Turning
 17  Milling
 18  Wood-to-Metal Analogy
 19  Quick Reference — Ferrous
 20  Quick Reference — Non-Ferrous
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

IMG_DIR = os.path.join(os.path.dirname(__file__), "img")

# ── Theme colors ──────────────────────────────────────────────
DARK_BG = RGBColor(0x1A, 0x1A, 0x2E)
STEEL_BLUE = RGBColor(0x4A, 0x6F, 0xA5)
YELLOW = RGBColor(0xE2, 0xB7, 0x14)
WHITE = RGBColor(0xF0, 0xF0, 0xF0)
LIGHT_GREY = RGBColor(0xC0, 0xC0, 0xC0)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ── Helper functions ──────────────────────────────────────────
def set_slide_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items,
                    font_size=16, color=WHITE, spacing_after=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = spacing_after
        run = p.add_run()
        run.text = f"  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Segoe UI"
        p.level = 0
    return txBox


def add_image_safe(slide, img_name, left, top, width=None, height=None):
    path = os.path.join(IMG_DIR, img_name)
    if os.path.exists(path):
        kwargs = {"image_file": path, "left": left, "top": top}
        if width:
            kwargs["width"] = width
        if height:
            kwargs["height"] = height
        return slide.shapes.add_picture(**kwargs)
    return None


def add_heading(slide, text, top=Inches(0.3), font_size=36, subtitle=None):
    left = Inches(0.6)
    width = Inches(12)
    if subtitle:
        txBox = slide.shapes.add_textbox(left, top, width, Inches(0.8))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        run1 = p.add_run()
        run1.text = text + " "
        run1.font.size = Pt(font_size)
        run1.font.color.rgb = WHITE
        run1.font.bold = True
        run1.font.name = "Segoe UI"
        run2 = p.add_run()
        run2.text = subtitle
        run2.font.size = Pt(24)
        run2.font.color.rgb = LIGHT_GREY
        run2.font.name = "Segoe UI"
    else:
        txBox = slide.shapes.add_textbox(left, top, width, Inches(0.8))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = WHITE
        p.font.bold = True
        p.font.name = "Segoe UI"
    # Yellow accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top + Inches(0.7), Inches(3), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = YELLOW
    line.line.fill.background()
    return txBox


def add_speaker_notes(slide, text):
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = text


def make_metal_profile(title, subtitle, items, img1, img2, notes_text):
    """Create a standard metal-profile slide (properties left, 2 images right)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    if subtitle:
        add_heading(slide, title, subtitle=subtitle)
    else:
        add_heading(slide, title)
    add_bullet_list(slide, Inches(0.6), Inches(1.3), Inches(7), Inches(6),
                    items, font_size=15, spacing_after=Pt(4))
    add_image_safe(slide, img1, Inches(8), Inches(1.3),
                   width=Inches(4.8), height=Inches(2.8))
    add_image_safe(slide, img2, Inches(8), Inches(4.3),
                   width=Inches(4.8), height=Inches(2.8))
    add_speaker_notes(slide, notes_text)
    return slide


def make_process_slide(title, process_subtitle, items, img, notes_text):
    """Create a process slide (bullets left, 1 large image right)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_heading(slide, title)
    add_textbox(slide, Inches(0.6), Inches(0.85), Inches(6), Inches(0.4),
                process_subtitle, font_size=18, color=LIGHT_GREY)
    add_bullet_list(slide, Inches(0.6), Inches(1.5), Inches(7), Inches(5.5),
                    items, font_size=16, spacing_after=Pt(8))
    add_image_safe(slide, img, Inches(8), Inches(1.5),
                   width=Inches(4.8), height=Inches(5))
    add_speaker_notes(slide, notes_text)
    return slide


def style_cell(cell, text, is_header=False, is_highlight=False, row_idx=0):
    cell.text = text
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.size = Pt(14) if not is_header else Pt(15)
        paragraph.font.bold = is_header or is_highlight
        paragraph.font.name = "Segoe UI"
        paragraph.alignment = PP_ALIGN.CENTER
        if is_header:
            paragraph.font.color.rgb = WHITE
        elif is_highlight:
            paragraph.font.color.rgb = YELLOW
        else:
            paragraph.font.color.rgb = WHITE
    tcPr = cell._tc.get_or_add_tcPr()
    solidFill = tcPr.makeelement(qn('a:solidFill'), {})
    if is_header:
        srgb = solidFill.makeelement(qn('a:srgbClr'), {'val': '3A5A8A'})
    elif row_idx % 2 == 0:
        srgb = solidFill.makeelement(qn('a:srgbClr'), {'val': '222238'})
    else:
        srgb = solidFill.makeelement(qn('a:srgbClr'), {'val': '1A1A2E'})
    solidFill.append(srgb)
    tcPr.append(solidFill)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


# ══════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_image_safe(slide, "title-bg-converted.jpg", Inches(0), Inches(0),
               width=SLIDE_W, height=SLIDE_H)
overlay = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_W, SLIDE_H)
overlay.fill.solid()
overlay.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
spPr = overlay._element.spPr
sf = spPr.find(qn('a:solidFill'))
if sf is not None:
    srgb = sf.find(qn('a:srgbClr'))
    if srgb is not None:
        alpha = srgb.makeelement(qn('a:alpha'), {})
        alpha.set('val', '70000')
        srgb.append(alpha)
overlay.line.fill.background()
add_textbox(slide, Inches(1), Inches(2.2), Inches(11), Inches(1.5),
            "Metal Types", font_size=60, color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, Inches(1), Inches(3.8), Inches(11), Inches(1),
            "Metals & Manufacturing", font_size=28, color=YELLOW,
            alignment=PP_ALIGN.CENTER)
add_speaker_notes(slide,
    "WELCOME & OPENING (2 min)\n\n"
    "Welcome to our unit on metal types. Just like we learn to identify different species "
    "of wood in woodworking, in metals class we need to be able to identify and understand "
    "the properties of different metals — because choosing the wrong metal for a job can "
    "mean the difference between a project that lasts a lifetime and one that fails on day one.\n\n"
    "Key question to pose: \"If I handed you two pieces of metal that looked almost identical, "
    "how would you tell them apart?\" (Let students brainstorm — weight, magnetism, color, "
    "spark test). We'll come back to this.")

# ══════════════════════════════════════════════════════════════
# SLIDE 2 — WHAT ABOUT METAL?
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_textbox(slide, Inches(1), Inches(2.5), Inches(11), Inches(1.5),
            "What About Metal?", font_size=52, color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)
line = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(5), Inches(4), Inches(3), Pt(3))
line.fill.solid()
line.fill.fore_color.rgb = YELLOW
line.line.fill.background()
add_speaker_notes(slide,
    "TRANSITION (1 min)\n\n"
    "In woodworking, students learn to tell the difference between hardwoods and softwoods — "
    "oak vs. pine, maple vs. cedar. Each one has a color, a grain, a hardness, and a best use.\n\n"
    "Metals are the same way. Steel isn't just \"steel.\" Aluminum isn't just \"aluminum.\" "
    "There are families, grades, and alloys — and each one behaves differently when you "
    "cut it, weld it, bend it, or finish it.\n\n"
    "Today we're going to build your metal vocabulary so you can walk into any shop and "
    "know what you're working with.")

# ══════════════════════════════════════════════════════════════
# SLIDE 3 — FERROUS vs. NON-FERROUS
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_heading(slide, "Terms: Ferrous vs. Non-Ferrous")
col_left = Inches(0.6)
add_textbox(slide, col_left, Inches(1.2), Inches(5.5), Inches(0.5),
            "Ferrous", font_size=28, color=YELLOW, bold=True)
add_image_safe(slide, "ferrous-rust.jpg", col_left, Inches(1.8),
               width=Inches(5.5), height=Inches(2.2))
add_bullet_list(slide, col_left, Inches(4.1), Inches(5.5), Inches(3),
    ["Contains iron (Fe)", "Magnetic", "Prone to rust (oxidation)",
     "Generally stronger", "Examples: steel, cast iron, wrought iron"],
    font_size=15)
col_right = Inches(6.8)
add_textbox(slide, col_right, Inches(1.2), Inches(5.5), Inches(0.5),
            "Non-Ferrous", font_size=28, color=YELLOW, bold=True)
add_image_safe(slide, "nonferrous-copper.jpg", col_right, Inches(1.8),
               width=Inches(5.5), height=Inches(2.2))
add_bullet_list(slide, col_right, Inches(4.1), Inches(5.5), Inches(3),
    ["Contains no iron", "Not magnetic", "Naturally corrosion resistant",
     "Generally lighter", "Examples: aluminum, copper, brass, titanium"],
    font_size=15)
add_speaker_notes(slide,
    "KEY CONCEPTS (5 min)\n\n"
    "This is the single most important classification in metals. Write \"Ferrous\" and "
    "\"Non-Ferrous\" on the board.\n\n"
    "Ferrous comes from the Latin word ferrum meaning iron — that's why iron's chemical "
    "symbol is Fe. ANY metal that contains iron as its primary ingredient is ferrous.\n\n"
    "Quick demo idea: Grab a magnet and walk around the shop. Touch it to different metals. "
    "If it sticks = ferrous. If it doesn't = non-ferrous.\n\n"
    "Check for understanding: \"Is stainless steel ferrous or non-ferrous?\" (Trick question — "
    "it IS ferrous because it contains iron, but the chromium content prevents rust.)")

# ══════════════════════════════════════════════════════════════
# SLIDE 4 — WHY DO YOU CARE?
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_heading(slide, "Why Do You Care?")
add_bullet_list(slide, Inches(1), Inches(1.5), Inches(11), Inches(5),
    ["The harder the metal, the more it wears your tools",
     "\"Machinability\" = how easy it is to cut, drill, and shape",
     "Wrong metal + wrong process = ruined project",
     "Different metals require different welding wire, gas, and settings",
     "Cost varies wildly — mild steel is ~$1/lb, stainless can be ~$3-4/lb"],
    font_size=22, spacing_after=Pt(14))
add_speaker_notes(slide,
    "RELEVANCE CONNECTION (3 min)\n\n"
    "This is the \"so what?\" slide. Metal selection directly affects shop work.\n\n"
    "Analogies to woodworking:\n"
    "- Sanding oak vs. pine = machining stainless vs. aluminum.\n"
    "- Balsa for a shelf = aluminum where you need steel strength.\n\n"
    "Real consequences:\n"
    "- MIG weld aluminum with steel wire = contaminated weld, scrapped part.\n"
    "- Lathe speeds for steel on aluminum = clogged tool, bad finish.\n"
    "- Pick stainless when mild steel works = 3x material cost.\n\n"
    "Engagement prompt: \"Has anyone tried to drill through something and the drill bit "
    "just squealed and smoked? What happened?\"")

# ══════════════════════════════════════════════════════════════
# SLIDE 5 — METAL IDENTIFICATION METHODS  (NEW)
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_heading(slide, "Metal Identification Methods")

# Four-quadrant layout
quad_w = Inches(5.8)
quad_h = Inches(2.5)
quad_fs = 14

# Top-left — Magnet Test
add_textbox(slide, Inches(0.6), Inches(1.2), quad_w, Inches(0.4),
            "1. Magnet Test", font_size=22, color=YELLOW, bold=True)
add_bullet_list(slide, Inches(0.6), Inches(1.7), quad_w, quad_h,
    ["Sticks = ferrous (contains iron)",
     "Doesn't stick = non-ferrous (aluminum, copper, brass)",
     "Caution: some stainless is weakly magnetic",
     "Simplest field test — always start here"],
    font_size=quad_fs, spacing_after=Pt(3))

# Top-right — Spark Test
add_textbox(slide, Inches(6.8), Inches(1.2), quad_w, Inches(0.4),
            "2. Spark Test (Grinder)", font_size=22, color=YELLOW, bold=True)
add_bullet_list(slide, Inches(6.8), Inches(1.7), quad_w, quad_h,
    ["Mild steel: long white sparks, many forks (carbon bursts)",
     "Stainless: short dark-red sparks, very few forks",
     "Cast iron: short red sparks, very fine streams",
     "Aluminum & copper: no sparks at all"],
    font_size=quad_fs, spacing_after=Pt(3))

# Bottom-left — Weight / Density
add_textbox(slide, Inches(0.6), Inches(4.0), quad_w, Inches(0.4),
            "3. Weight / Density Feel", font_size=22, color=YELLOW, bold=True)
add_bullet_list(slide, Inches(0.6), Inches(4.5), quad_w, quad_h,
    ["Aluminum: shockingly light (~1/3 of steel)",
     "Steel & stainless: heavy, similar to each other",
     "Copper & brass: noticeably heavier than steel",
     "Titanium: lighter than steel, heavier than aluminum"],
    font_size=quad_fs, spacing_after=Pt(3))

# Bottom-right — Visual / Color
add_textbox(slide, Inches(6.8), Inches(4.0), quad_w, Inches(0.4),
            "4. Visual / Color Cues", font_size=22, color=YELLOW, bold=True)
add_bullet_list(slide, Inches(6.8), Inches(4.5), quad_w, quad_h,
    ["Mild steel: dark grey with bluish mill scale",
     "Aluminum: bright silver-white, stays clean",
     "Copper: distinctive reddish-orange, develops green patina",
     "Brass: gold/yellow tone — unmistakable"],
    font_size=quad_fs, spacing_after=Pt(3))

add_speaker_notes(slide,
    "METAL IDENTIFICATION (5 min)\n\n"
    "In woodworking, students ID species by grain, color, and smell. Metals use different "
    "methods — but the same principle: systematic testing.\n\n"
    "MAGNET TEST — always start here. It's the fastest sort into ferrous vs. non-ferrous. "
    "Demo: walk around the shop with a magnet.\n\n"
    "SPARK TEST — touch metal to bench grinder and observe the spark pattern. This is how "
    "experienced machinists ID unknown steel. Safety: wear face shield, keep fingers clear.\n"
    "- Mild steel: long, white, bushy sparks with lots of carbon forks\n"
    "- Stainless: short, dark-red, few forks (low carbon)\n"
    "- Cast iron: short, red, fine streams\n"
    "- Non-ferrous metals don't spark at all\n\n"
    "WEIGHT TEST — have students hold same-size pieces of different metals. The difference "
    "between aluminum and steel is dramatic.\n\n"
    "VISUAL TEST — color is less reliable than the others (freshly ground surfaces look "
    "similar) but mill scale, patina, and surface finish give clues.\n\n"
    "Activity idea: Set up a station with 5-6 unmarked metal samples. Students use all four "
    "tests to identify each one. Give them a worksheet to record results.")

# ══════════════════════════════════════════════════════════════
# SLIDE 6 — MILD STEEL (1018 / A36)
# ══════════════════════════════════════════════════════════════
make_metal_profile(
    "Mild Steel", "(1018 / A36)",
    ["FERROUS",
     "Color: dark grey, bluish mill scale",
     "Hardness: low-medium (Brinell ~126)",
     "Machinability: excellent",
     "Weldability: excellent — easiest to weld",
     "Density: 7.87 g/cm\u00B3 (heavy)",
     "Magnetic: yes",
     "Corrosion: rusts easily — oil or prime it",
     "Carbon content: 0.18% (low carbon)",
     "Cost: ~$1/lb — the cheapest structural metal",
     "The \"pine\" of metals — learn on this first"],
    "mild-steel.jpg", "mild-steel-project.jpg",
    "MILD STEEL DEEP DIVE (5 min)\n\n"
    "Mild steel is the bread and butter of any metals shop.\n\n"
    "Naming: \"1018\" = AISI/SAE grade. \"10\" = plain carbon, \"18\" = 0.18% carbon. "
    "\"A36\" = ASTM structural grade. Most flat bar, angle, and tube in shop is A36.\n\n"
    "Identification:\n"
    "- Silvery grey with dark blue-grey mill scale\n"
    "- Magnetic — always sticks\n"
    "- Spark test: long, white, branching with lots of forks\n"
    "- WILL rust bare. Oil or prime immediately.\n\n"
    "In our shop: MIG weld with ER70S-6 wire, C25 gas. Machines beautifully.\n\n"
    "Analogy: \"Mild steel is to metals what pine is to woodworking.\"")

# ══════════════════════════════════════════════════════════════
# SLIDE 7 — 4140 CHROME-MOLY STEEL  (NEW)
# ══════════════════════════════════════════════════════════════
make_metal_profile(
    "4140 Chrome-Moly Steel", "(4140)",
    ["FERROUS",
     "Color: dark grey, similar to mild steel",
     "Hardness: medium-high (Brinell ~197 annealed, ~500+ hardened)",
     "Machinability: good when annealed, poor when hardened",
     "Weldability: weldable — preheat required (400\u00B0F)",
     "Density: 7.85 g/cm\u00B3 (same as mild steel)",
     "Magnetic: yes",
     "Corrosion: rusts — slightly better than mild steel",
     "Alloy: 1% chromium + 0.2% molybdenum + 0.4% carbon",
     "Heat treatable — can be hardened & tempered",
     "The \"red oak\" of metals — harder, more demanding, premium result"],
    "mild-steel.jpg", "mild-steel-project.jpg",
    "4140 CHROME-MOLY DEEP DIVE (5 min)\n\n"
    "4140 is the step up from mild steel. When a project needs more strength or wear "
    "resistance, this is what you reach for.\n\n"
    "Naming: \"41\" = chromium-molybdenum alloy family, \"40\" = 0.40% carbon. More carbon "
    "= harder, less ductile. The chromium and moly add toughness and hardenability.\n\n"
    "Key difference from mild steel: 4140 can be HEAT TREATED. You can harden it to "
    "Rockwell 55+ then temper it back to the toughness you need. Mild steel can't do this.\n\n"
    "In our shop:\n"
    "- Machines well in the annealed state (buy it annealed)\n"
    "- Welding: need to preheat to ~400F and slow-cool to prevent cracking\n"
    "- Used for: axles, shafts, gears, tooling, motorsport parts, roll cage tubes\n"
    "- 4130 (similar alloy, less carbon) is the go-to for roll cages and bike frames\n\n"
    "Analogy: \"4140 is the red oak of metals — you pay more, it's harder to work, but "
    "the finished product is a cut above.\"")

# ══════════════════════════════════════════════════════════════
# SLIDE 8 — TOOL STEEL (A2 / D2)  (NEW)
# ══════════════════════════════════════════════════════════════
make_metal_profile(
    "Tool Steel", "(A2 / D2)",
    ["FERROUS",
     "Color: dark grey — looks like other steels until hardened",
     "Hardness: very high (Rockwell 58-62 HRC hardened)",
     "Machinability: fair annealed, extremely poor hardened",
     "Weldability: very difficult — risk of cracking",
     "Density: 7.86 g/cm\u00B3",
     "Magnetic: yes",
     "Corrosion: moderate (D2 has 12% chromium — semi-stainless)",
     "Carbon: 1.0% (A2) to 1.5% (D2) — very high carbon",
     "Used for: dies, punches, cutting tools, blades, jigs",
     "The \"exotic hardwood\" of metals — specialized, unforgiving"],
    "mild-steel.jpg", "mild-steel-project.jpg",
    "TOOL STEEL DEEP DIVE (5 min)\n\n"
    "Tool steel is what we make our tools FROM. When you need a cutting edge that stays "
    "sharp or a die that stamps thousands of parts without deforming.\n\n"
    "Naming:\n"
    "- A2 = air-hardening tool steel. Heat it, let it air cool — it hardens on its own.\n"
    "- D2 = high-carbon, high-chromium (12%) die steel. Sometimes called semi-stainless.\n"
    "- O1 = oil-hardening tool steel — quench in oil. Traditional knifemaker's steel.\n\n"
    "In our shop:\n"
    "- Students won't machine hardened tool steel (it would destroy our tooling)\n"
    "- You CAN machine it in the annealed state, then send it out for heat treatment\n"
    "- Used for making custom fixtures, punches, stamps, and project components that need wear resistance\n\n"
    "Safety: hardened tool steel is BRITTLE. If you drop it, it can crack or chip. "
    "Flying chips from tool steel are extremely sharp.\n\n"
    "Analogy: \"Tool steel is like ebony or cocobolo — you don't use it for everything, "
    "but when you need extreme hardness and edge retention, nothing else will do.\"")

# ══════════════════════════════════════════════════════════════
# SLIDE 9 — CAST IRON  (NEW)
# ══════════════════════════════════════════════════════════════
make_metal_profile(
    "Cast Iron", "(Grey / Ductile)",
    ["FERROUS",
     "Color: dark grey, rough sandy surface (as-cast)",
     "Hardness: medium-high (Brinell ~200 grey, ~250 ductile)",
     "Machinability: excellent (grey iron) — chips are powder/dust",
     "Weldability: very difficult — preheating and special rods required",
     "Density: 7.15 g/cm\u00B3 (slightly lighter than steel)",
     "Magnetic: yes",
     "Corrosion: rusts, but slower than mild steel",
     "Carbon: 2-4% — much higher than any steel",
     "Brittle — does NOT bend, it snaps",
     "The \"workhorse\" of machine frames — every lathe and mill sits on it"],
    "mild-steel.jpg", "mild-steel-project.jpg",
    "CAST IRON DEEP DIVE (5 min)\n\n"
    "Cast iron is everywhere in the shop — but not as raw material for student projects. "
    "It's what your machines are MADE of.\n\n"
    "Why teach it: students need to understand that the lathe bed, the mill table, the "
    "vise body — these are all cast iron. It absorbs vibration (dampening), it's extremely "
    "rigid, and it machines to a beautiful flat surface.\n\n"
    "Types:\n"
    "- Grey iron: most common. The graphite flakes make it easy to machine (chips come off "
    "as powder, not curls). But it's brittle — don't drop a vise!\n"
    "- Ductile iron: graphite is spheroidal, giving it some flex. Used for pipe fittings.\n\n"
    "Identification:\n"
    "- Very dark grey, rough/sandy texture on unfinished surfaces\n"
    "- Machined surfaces are smooth, dark grey with a slight sheen\n"
    "- Spark test: short, red, fine streams — less showy than steel\n"
    "- Chips are powdery (grey iron) — be careful, they stain everything\n\n"
    "Shop rule: NEVER weld cast iron without talking to the instructor first. It requires "
    "preheat to 500F+ and special nickel rods. A cold weld WILL crack.\n\n"
    "Where you see it: engine blocks, brake rotors, manhole covers, pipe fittings, "
    "machine tool frames, woodworking bench vises.")

# ══════════════════════════════════════════════════════════════
# SLIDE 10 — 6061 ALUMINUM
# ══════════════════════════════════════════════════════════════
make_metal_profile(
    "6061 Aluminum", None,
    ["NON-FERROUS",
     "Color: silver-white, bright",
     "Hardness: medium (Brinell ~95)",
     "Machinability: very good",
     "Weldability: weldable (TIG preferred, MIG possible)",
     "Density: 2.70 g/cm\u00B3 (~1/3 the weight of steel)",
     "Magnetic: no",
     "Corrosion: naturally resistant (oxide layer)",
     "Alloy series: 6xxx (magnesium + silicon)",
     "Heat treatable — T6 temper is strongest",
     "The \"poplar\" of metals — light, versatile, everywhere"],
    "aluminum.jpg", "aluminum-project.jpg",
    "6061 ALUMINUM DEEP DIVE (5 min)\n\n"
    "If mild steel is the pine, 6061 aluminum is the poplar — lighter, softer, very versatile.\n\n"
    "Naming: \"6061\" — the \"6\" means magnesium + silicon. \"6061-T6\" = heat treated for max strength.\n\n"
    "Identification:\n"
    "- Bright silver-white, much lighter looking than steel\n"
    "- Pick it up — shockingly light. 1/3 the weight of steel.\n"
    "- NOT magnetic. No spark on grinder.\n\n"
    "In our shop:\n"
    "- Machines FAST — 2-3x spindle speed vs. steel\n"
    "- \"Gummy\" — chips can weld to cutting tool without proper speed/feed\n"
    "- Welding: aluminum wire, 100% argon, AC (TIG) or special MIG setup\n\n"
    "Where you see it: Bike frames, aircraft, phones, laptops, car wheels.")

# ══════════════════════════════════════════════════════════════
# SLIDE 11 — COPPER  (NEW)
# ══════════════════════════════════════════════════════════════
make_metal_profile(
    "Copper", "(C110 / C101)",
    ["NON-FERROUS",
     "Color: distinctive reddish-orange (unmistakable)",
     "Hardness: low (Brinell ~50) — very soft",
     "Machinability: fair — gummy and smears",
     "Weldability: solderable & brazeable, TIG weldable",
     "Density: 8.96 g/cm\u00B3 (heavier than steel!)",
     "Magnetic: no",
     "Corrosion: develops green patina (verdigris), never rusts",
     "Conductivity: best of any common metal (electrical & thermal)",
     "Cost: ~$4-5/lb — expensive, but scrap value is high",
     "The \"cherry\" of metals — beautiful, functional, mid-range difficulty"],
    "nonferrous-copper.jpg", "nonferrous-copper.jpg",
    "COPPER DEEP DIVE (5 min)\n\n"
    "Copper is one of the oldest metals humans have worked — the Copper Age predates the "
    "Bronze Age. Students often recognize it from plumbing and electrical work.\n\n"
    "Why it matters:\n"
    "- #1 electrical conductor (after silver). Every wire in this building is copper.\n"
    "- #1 thermal conductor in common use. That's why pots and pans have copper bottoms.\n"
    "- Antimicrobial — copper surfaces kill bacteria. Hospital door handles.\n\n"
    "Identification:\n"
    "- Color is unmistakable — reddish-orange, develops green patina over time\n"
    "- Noticeably HEAVY — denser than steel\n"
    "- Not magnetic\n"
    "- Very soft — you can scratch it with a fingernail (pure copper)\n\n"
    "In our shop:\n"
    "- Used for electrical projects, decorative work, art pieces\n"
    "- Soldering and brazing are more common than welding\n"
    "- Machining: tends to be gummy. Sharp tools, moderate speed, use cutting fluid.\n"
    "- Scrap copper has real value — NEVER throw it away\n\n"
    "Where you see it: Electrical wiring, plumbing pipes, roofing (Statue of Liberty!), "
    "heat exchangers, decorative art, circuit boards.")

# ══════════════════════════════════════════════════════════════
# SLIDE 12 — BRASS  (NEW)
# ══════════════════════════════════════════════════════════════
make_metal_profile(
    "Brass", "(360 Free-Machining / 260 Cartridge)",
    ["NON-FERROUS",
     "Color: gold/yellow — unmistakable",
     "Hardness: low-medium (Brinell ~120 for 360 brass)",
     "Machinability: outstanding — the best machining metal",
     "Weldability: solderable & brazeable easily",
     "Density: 8.50 g/cm\u00B3 (heavier than steel)",
     "Magnetic: no",
     "Corrosion: excellent resistance, does not rust",
     "Composition: copper + zinc alloy (60-70% Cu, 30-40% Zn)",
     "Non-sparking — safe for explosive environments",
     "The \"maple\" of metals — machines beautifully, premium look"],
    "nonferrous-copper.jpg", "nonferrous-copper.jpg",
    "BRASS DEEP DIVE (5 min)\n\n"
    "Brass is a machinist's dream. If you want students to experience what a perfect "
    "finish looks like off the lathe, hand them a piece of 360 brass.\n\n"
    "Types:\n"
    "- 360 (free-machining brass): 3% lead added for chip breaking. Creates beautiful, "
    "short chips that fly away clean. The gold standard for screw machine parts.\n"
    "- 260 (cartridge brass): 70/30 copper/zinc. More ductile, used for ammunition, "
    "musical instruments, decorative hardware.\n\n"
    "Identification:\n"
    "- Gold/yellow color — impossible to confuse with anything else\n"
    "- Heavy — noticeably heavier than steel for the same size\n"
    "- Not magnetic\n"
    "- Rings like a bell when you tap it (compared to steel's dull clang)\n\n"
    "In our shop:\n"
    "- Lathe projects in brass look incredible with minimal polishing\n"
    "- Great for decorative hardware, handles, fittings, nameplates\n"
    "- NON-SPARKING — this is why tools for use around explosives are made from brass\n"
    "- Scrap brass has good value — keep it separate from steel and aluminum\n\n"
    "Where you see it: Door hardware, musical instruments (trumpets!), ammunition casings, "
    "plumbing fittings, decorative trim, marine hardware.")

# ══════════════════════════════════════════════════════════════
# SLIDE 13 — TITANIUM  (NEW)
# ══════════════════════════════════════════════════════════════
make_metal_profile(
    "Titanium", "(Grade 5 / Ti-6Al-4V)",
    ["NON-FERROUS",
     "Color: medium grey with slight blue/purple tint",
     "Hardness: high (Brinell ~334 for Grade 5)",
     "Machinability: poor — generates extreme heat",
     "Weldability: TIG only, requires argon shielding everywhere",
     "Density: 4.43 g/cm\u00B3 (lighter than steel, heavier than aluminum)",
     "Magnetic: no",
     "Corrosion: exceptional — nearly immune",
     "Strength-to-weight: best of any common metal",
     "Cost: ~$15-25/lb — extremely expensive",
     "The \"purple heart\" of metals — exotic, hard, impressive"],
    "stainless-steel.jpg", "stainless-project.jpg",
    "TITANIUM DEEP DIVE (5 min)\n\n"
    "Titanium is the aerospace wonder metal. Students may not work with it in our shop, "
    "but they need to know it exists and why it's special.\n\n"
    "Naming: Grade 5 (Ti-6Al-4V) = 6% aluminum, 4% vanadium. This is the most common "
    "titanium alloy — used in 50%+ of all titanium applications.\n\n"
    "Why it's special:\n"
    "- As strong as steel at 60% of the weight\n"
    "- As corrosion-resistant as stainless (better, actually)\n"
    "- Biocompatible — your body won't reject it (hip replacements, dental implants)\n"
    "- Heat resistant — jet engines operate at 1000F+\n\n"
    "Why it's difficult:\n"
    "- Generates extreme heat when machined — requires flood coolant\n"
    "- Work hardens like stainless but worse\n"
    "- Welding MUST be shielded on BOTH sides — any oxygen contamination turns it brittle\n"
    "- Expensive: a 1\" round bar of Ti-6Al-4V costs what a full stick of mild steel costs\n\n"
    "Identification:\n"
    "- Lighter than steel but heavier than aluminum (distinct \"in between\" feel)\n"
    "- Not magnetic\n"
    "- Spark test: brilliant white sparks (distinctive)\n"
    "- When ground, can develop blue/purple oxide colors\n\n"
    "Where you see it: Aircraft frames, jet engines, surgical implants, high-end bicycle "
    "frames, spacecraft, premium watches, golf clubs.")

# ══════════════════════════════════════════════════════════════
# SLIDE 14 — STAINLESS STEEL (304 / 316)
# ══════════════════════════════════════════════════════════════
make_metal_profile(
    "Stainless Steel", "(304 / 316)",
    ["FERROUS (still contains iron!)",
     "Color: bright silver, stays shiny",
     "Hardness: high (Brinell ~200)",
     "Machinability: poor — work hardens",
     "Weldability: weldable (TIG preferred)",
     "Density: 8.00 g/cm\u00B3 (slightly heavier than mild steel)",
     "Magnetic: depends on grade (304 weakly magnetic)",
     "Corrosion: highly resistant (10.5%+ chromium)",
     "Cost: ~$3-4/lb — 3-4x mild steel",
     "Don't cross-contaminate with carbon steel tools!",
     "The \"premium\" metal — worth it when corrosion matters"],
    "stainless-steel.jpg", "stainless-project.jpg",
    "STAINLESS STEEL DEEP DIVE (5 min)\n\n"
    "Stainless steel is the most misunderstood metal in the shop.\n\n"
    "Naming:\n"
    "- 304: most common. 18% chromium, 8% nickel (\"18-8\").\n"
    "- 316: added molybdenum for marine/medical.\n"
    "- Chromium (min 10.5%) forms invisible, self-healing protective layer.\n\n"
    "Identification:\n"
    "- Stays bright — no mill scale, no rust.\n"
    "- Slightly heavier than mild steel.\n"
    "- Magnet test unreliable — some grades are slightly magnetic.\n"
    "- Spark test: short, dark-red, very few forks.\n\n"
    "Critical shop rules:\n"
    "- NEVER use tools on stainless that touched carbon steel (cross-contamination = rust spots)\n"
    "- Keep dedicated stainless tools marked separately\n"
    "- Work hardens — keep feeds aggressive, tools sharp\n\n"
    "Where you see it: Kitchen, food processing, surgical, marine, architecture.")

# ══════════════════════════════════════════════════════════════
# SLIDE 15 — MIG WELDING
# ══════════════════════════════════════════════════════════════
make_process_slide("MIG Welding", "Gas Metal Arc Welding (GMAW)",
    ["How it works: wire-fed electrode melts into the joint",
     "Shielding gas protects the weld pool from contamination",
     "Steel setup: ER70S-6 wire + C25 gas (75% Ar / 25% CO2)",
     "Aluminum setup: ER4043 wire + 100% Argon",
     "Easiest welding process to learn",
     "\"Point and shoot\" — the hot glue gun of welding",
     "Travel speed matters — too fast = cold weld, too slow = burn-through",
     "Wire speed + voltage must be balanced"],
    "mig-welding.jpg",
    "MIG WELDING OVERVIEW (5 min)\n\n"
    "MIG is the first welding process most students learn.\n\n"
    "Explain: \"Hot glue gun, but feeds metal wire. Arc is ~6,000F. "
    "Shielding gas is an invisible bubble keeping oxygen out.\"\n\n"
    "Key settings:\n"
    "- Wire speed: filler amount (too much = blobby, too little = no fusion)\n"
    "- Voltage: arc length and heat\n"
    "- Travel speed: hardest for beginners\n\n"
    "Common mistakes: gun too far, moving too fast, dirty metal, wrong polarity.\n\n"
    "Safety: Never look at arc without helmet (min shade 10). Arc eye is no joke.")

# ══════════════════════════════════════════════════════════════
# SLIDE 16 — LATHE TURNING
# ══════════════════════════════════════════════════════════════
make_process_slide("Lathe Turning", "Cylindrical Machining",
    ["How it works: workpiece spins, cutting tool stays still",
     "Creates round/cylindrical shapes",
     "Operations: facing, turning, boring, threading, knurling",
     "RPM changes by material:",
     "    Steel: ~300-600 RPM",
     "    Aluminum: ~800-1500 RPM",
     "    Stainless: ~200-400 RPM (slow!)",
     "Cutting fluid reduces heat and improves finish",
     "Never leave the chuck key in the chuck"],
    "lathe.jpg",
    "LATHE TURNING OVERVIEW (5 min)\n\n"
    "The lathe: oldest machine tool, still essential.\n\n"
    "Key concept: WORK spins, TOOL stays still. Opposite of drill/mill. "
    "Because work spins, everything a lathe makes is round.\n\n"
    "Speed:\n"
    "- Aluminum: high RPMs (soft, conducts heat)\n"
    "- Stainless: LOW RPMs (work-hardens)\n"
    "- Rule: softer = faster RPM, harder = slower\n\n"
    "Safety:\n"
    "1. NEVER leave chuck key in chuck — projectile\n"
    "2. No loose clothing, gloves, or jewelry\n"
    "3. Safety glasses always\n"
    "4. Stand to side of chuck, not in front")

# ══════════════════════════════════════════════════════════════
# SLIDE 17 — MILLING
# ══════════════════════════════════════════════════════════════
make_process_slide("Milling", "Precision Material Removal",
    ["How it works: spinning cutter removes material from clamped workpiece",
     "Creates flat surfaces, slots, pockets, and contours",
     "Operations: face milling, end milling, slot cutting, drilling",
     "Feed rate changes by material:",
     "    Steel: moderate feed, lower RPM",
     "    Aluminum: aggressive feed, higher RPM",
     "    Stainless: light feed, low RPM, constant cutting fluid",
     "Workholding is critical — vise must be tight",
     "Climb milling vs. conventional — know the difference"],
    "milling.jpg",
    "MILLING OVERVIEW (5 min)\n\n"
    "The mill makes everything the lathe can't — flat surfaces, square shapes, pockets.\n\n"
    "Key concept: TOOL spins, WORK stays still (on movable table). Multiple flutes, "
    "each takes a small bite.\n\n"
    "Material tips:\n"
    "- Aluminum: 2-3 flute end mills (chip evacuation)\n"
    "- Steel: 4 flute, slower RPM\n"
    "- Stainless: sharp tools critical (dull = work-hardening cycle)\n\n"
    "Safety:\n"
    "- Workpiece MUST be securely clamped\n"
    "- Chips with brush, never hands\n"
    "- Safety glasses + face shield for heavy cuts")

# ══════════════════════════════════════════════════════════════
# SLIDE 18 — WOOD-TO-METAL ANALOGY  (NEW)
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_heading(slide, "Wood-to-Metal Analogy")
add_textbox(slide, Inches(0.6), Inches(0.85), Inches(10), Inches(0.4),
            "Bridging what you know to what you\u2019re learning",
            font_size=18, color=LIGHT_GREY)

# Build the analogy table
rows, cols = 9, 4
tbl_left = Inches(0.8)
tbl_top = Inches(1.4)
tbl_width = Inches(11.7)
tbl_height = Inches(5.6)
table_shape = slide.shapes.add_table(rows, cols, tbl_left, tbl_top,
                                      tbl_width, tbl_height)
table = table_shape.table
for i, w in enumerate([Inches(2.5), Inches(3.2), Inches(3.0), Inches(3.0)]):
    table.columns[i].width = w

analogy_headers = ["Wood", "Metal Equivalent", "Why the Match", "Difficulty"]
analogy_data = [
    ["Pine",        "Mild Steel",       "Cheap, common, easy to work, learn on this",       "Beginner"],
    ["Poplar",      "6061 Aluminum",    "Light, versatile, softer, everywhere",             "Beginner"],
    ["Red Oak",     "4140 Cr-Mo Steel", "Harder, more demanding, premium result",           "Intermediate"],
    ["Maple",       "Brass",            "Machines beautifully, polishes to a premium look",  "Intermediate"],
    ["Cherry",      "Copper",           "Beautiful finish, mid-range difficulty, ages well",  "Intermediate"],
    ["Walnut",      "Stainless Steel",  "Premium, harder to work, worth the effort",         "Advanced"],
    ["Purple Heart", "Titanium",        "Exotic, expensive, hard, impressive result",        "Expert"],
    ["White Oak",   "Cast Iron",        "Heavy, durable, structural, been used for centuries", "Specialized"],
]

for j, h in enumerate(analogy_headers):
    style_cell(table.cell(0, j), h, is_header=True)
for i, row_data in enumerate(analogy_data):
    for j, val in enumerate(row_data):
        is_hl = (j == 1)  # Highlight the metal column
        style_cell(table.cell(i + 1, j), val, is_highlight=is_hl, row_idx=i + 1)

add_speaker_notes(slide,
    "WOOD-TO-METAL ANALOGY (5 min)\n\n"
    "This slide bridges what students already know from woodworking to the new metal "
    "vocabulary they're building. Go through each row and ask if anyone has worked with "
    "that wood species — then explain how the metal equivalent behaves similarly.\n\n"
    "Key teaching points:\n"
    "- Pine/Mild Steel: \"You learn on pine because it's forgiving. Same with mild steel. "
    "If you mess up, the material cost is low and the process is easy.\"\n"
    "- Poplar/Aluminum: \"Both are light, both are versatile, both are softer than their "
    "premium cousins.\"\n"
    "- Red Oak/4140: \"Red oak demands sharp tools and proper technique. So does 4140. But "
    "the finished product is noticeably better.\"\n"
    "- Cherry/Copper: \"Both develop a beautiful patina over time. Both require care to "
    "get a great finish.\"\n"
    "- Purple Heart/Titanium: \"Both are exotic, expensive, and hard — but when someone "
    "pulls off a project in either material, it's impressive.\"\n"
    "- White Oak/Cast Iron: \"Both have been used for structural applications for centuries. "
    "Heavy, durable, and not going anywhere.\"\n\n"
    "The difficulty column reinforces the idea of progression: start with beginner "
    "materials, build skills, then tackle harder metals as confidence grows.")

# ══════════════════════════════════════════════════════════════
# SLIDE 19 — QUICK REFERENCE: FERROUS METALS  (EXPANDED)
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_heading(slide, "Quick Reference", subtitle="Ferrous Metals")

rows, cols = 8, 6
tbl_left = Inches(0.4)
tbl_top = Inches(1.4)
tbl_width = Inches(12.5)
tbl_height = Inches(5.5)
table_shape = slide.shapes.add_table(rows, cols, tbl_left, tbl_top,
                                      tbl_width, tbl_height)
table = table_shape.table
for i, w in enumerate([Inches(1.8), Inches(2.1), Inches(2.1), Inches(2.1), Inches(2.2), Inches(2.2)]):
    table.columns[i].width = w

ferrous_headers = ["Property", "Mild Steel", "4140 Cr-Mo", "Tool Steel", "Cast Iron", "Stainless"]
ferrous_data = [
    ["Hardness",      "Low-Med",  "Med-High", "Very High", "Med-High", "High"],
    ["Machinability",  "Excellent", "Good",    "Fair",      "Excellent", "Poor"],
    ["Weldability",    "Excellent", "Moderate", "Difficult", "Difficult", "Moderate"],
    ["Rusts?",         "Yes!",     "Yes",      "Moderate",  "Slowly",    "No"],
    ["Heat Treatable", "No",       "Yes",      "Yes",       "No",        "No"],
    ["Magnetic?",      "Yes",      "Yes",      "Yes",       "Yes",       "Varies"],
    ["Relative Cost",  "$",        "$$",       "$$$",       "$",         "$$$"],
]
ferrous_highlights = {
    (1, 1): True,  # Excellent machinability mild steel
    (2, 1): True,  # Excellent weldability mild steel
    (3, 1): True,  # Yes! rusts
    (2, 4): True,  # Excellent machinability cast iron
    (5, 1): True,  # Yes heat treatable 4140
    (5, 2): True,  # Yes heat treatable tool steel
    (4, 5): True,  # No rust stainless
    (7, 4): True,  # $ cast iron cheap
    (7, 5): True,  # $$$ stainless expensive
}

for j, h in enumerate(ferrous_headers):
    style_cell(table.cell(0, j), h, is_header=True)
for i, row_data in enumerate(ferrous_data):
    for j, val in enumerate(row_data):
        is_hl = (i + 1, j) in ferrous_highlights
        style_cell(table.cell(i + 1, j), val, is_highlight=is_hl, row_idx=i + 1)

add_speaker_notes(slide,
    "FERROUS QUICK REFERENCE (3 min)\n\n"
    "Walk through the table column by column. Key talking points:\n"
    "- Mild Steel: best all-around for beginners. Cheap, welds great, machines great.\n"
    "- 4140: when you need more strength + heat treatability. Motorsport, tooling.\n"
    "- Tool Steel: extreme hardness. Dies, punches, blades. Don't machine it hardened.\n"
    "- Cast Iron: excellent machining (powder chips!), but don't try to weld it.\n"
    "- Stainless: corrosion king, but difficult to machine and expensive.\n\n"
    "Activity: \"If you're building a go-kart frame on a budget, which column do you pick?\" "
    "(Mild steel.) \"If the frame needs to be heat treated for a racing class?\" (4140.)")

# ══════════════════════════════════════════════════════════════
# SLIDE 20 — QUICK REFERENCE: NON-FERROUS METALS  (EXPANDED)
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide)
add_heading(slide, "Quick Reference", subtitle="Non-Ferrous Metals")

rows, cols = 8, 5
tbl_left = Inches(0.6)
tbl_top = Inches(1.4)
tbl_width = Inches(12)
tbl_height = Inches(5.5)
table_shape = slide.shapes.add_table(rows, cols, tbl_left, tbl_top,
                                      tbl_width, tbl_height)
table = table_shape.table
for i, w in enumerate([Inches(2.0), Inches(2.5), Inches(2.5), Inches(2.5), Inches(2.5)]):
    table.columns[i].width = w

nonferrous_headers = ["Property", "6061 Aluminum", "Copper", "Brass", "Titanium"]
nonferrous_data = [
    ["Weight",        "Light (~1/3)",  "Heavy",        "Heavy",         "Medium"],
    ["Machinability", "Very Good",     "Fair",         "Outstanding",   "Poor"],
    ["Weldability",   "Moderate",      "Solderable",   "Solderable",    "TIG only"],
    ["Corrosion",     "Resistant",     "Patina",       "Resistant",     "Exceptional"],
    ["Magnetic?",     "No",            "No",           "No",            "No"],
    ["Special",       "1/3 weight",    "Best conductor", "Non-sparking", "Strength/weight"],
    ["Relative Cost", "$$",            "$$$",          "$$",            "$$$$"],
]
nonferrous_highlights = {
    (1, 1): True,  # Light aluminum
    (2, 3): True,  # Outstanding brass machinability
    (4, 4): True,  # Exceptional titanium corrosion
    (6, 1): True,  # 1/3 weight aluminum
    (6, 2): True,  # Best conductor copper
    (6, 3): True,  # Non-sparking brass
    (6, 4): True,  # Strength/weight titanium
    (7, 4): True,  # $$$$ titanium
}

for j, h in enumerate(nonferrous_headers):
    style_cell(table.cell(0, j), h, is_header=True)
for i, row_data in enumerate(nonferrous_data):
    for j, val in enumerate(row_data):
        is_hl = (i + 1, j) in nonferrous_highlights
        style_cell(table.cell(i + 1, j), val, is_highlight=is_hl, row_idx=i + 1)

add_speaker_notes(slide,
    "NON-FERROUS QUICK REFERENCE (3 min)\n\n"
    "Key talking points per column:\n"
    "- 6061 Aluminum: the lightweight workhorse. When weight matters.\n"
    "- Copper: electrical/thermal champion. Beautiful but heavy and gummy.\n"
    "- Brass: the machinist's dream metal. Makes everything look premium.\n"
    "- Titanium: the exotic — strongest per pound, most corrosion-resistant, most expensive.\n\n"
    "Notice ALL non-ferrous metals are non-magnetic. That's your first sort.\n\n"
    "Discussion: \"You're designing a lightweight drone bracket that needs to survive "
    "outdoors. Which column?\" (Aluminum for cost/weight. Titanium if budget is unlimited.)\n\n"
    "Key takeaway: There's no \"best\" metal. Only the best metal FOR THE JOB. Every "
    "choice is a trade-off between strength, weight, cost, corrosion resistance, and workability.")

# ══════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════
output_path = os.path.join(os.path.dirname(__file__), "Metal_Types_Presentation.pptx")
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
